import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs, title):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = "Introduction to AI Course"

    return slide


def create_section_slide(prs, title):
    """Create a section title slide"""
    slide_layout = prs.slide_layouts[1]  # Section Header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    return slide


def create_content_slide(prs, title, content, is_code=False):
    """Create a content slide with bullet points or code"""
    slide_layout = prs.slide_layouts[1] if not is_code else prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if is_code:
        # For code slides
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        text_frame = code_box.text_frame
        p = text_frame.add_paragraph()
        p.text = content
        p.font.name = "Courier New"
        p.font.size = Pt(11)
    else:
        # For bullet point slides
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame

        for item in content:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

    return slide


def create_image_slide(prs, title, image_path, caption=None):
    """Create a slide with an image"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Add image if the file exists
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(2), Inches(2), width=Inches(6))

    # Add caption if provided
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        text_frame = caption_box.text_frame
        p = text_frame.add_paragraph()
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        p.font.italic = True

    return slide


def create_table_slide(prs, title, data):
    """Create a slide with a table"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Define table dimensions
    rows = len(data)
    cols = len(data[0]) if rows > 0 else 0

    if rows > 0 and cols > 0:
        # Calculate table position and size
        top = Inches(1.5)
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(4)

        # Add table to slide
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Populate table
        for row_idx, row in enumerate(data):
            for col_idx, cell_text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = cell_text

                # Make header row bold
                if row_idx == 0:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True

    return slide


def notebook_to_ppt(notebook_path, output_path):
    """Convert Jupyter notebook to PowerPoint presentation"""
    # Load notebook content
    with open(notebook_path, "r", encoding="utf-8") as f:
        notebook_content = json.load(f)

    # Create presentation
    prs = Presentation()

    # Parse cells and create slides
    title = "Introduction to Artificial Intelligence (AI)"
    create_title_slide(prs, title)

    for cell in notebook_content["cells"]:
        if cell["cell_type"] == "markdown":
            source = "".join(cell["source"])

            # Check if it's a section header (starts with ## or # )
            if source.startswith("## "):
                section_title = source.split("## ")[1].split("\n")[0]
                create_section_slide(prs, section_title)

                # Extract bullet points if any
                bullets = []
                for line in source.split("\n"):
                    if line.startswith("- "):
                        bullets.append(line[2:])

                if bullets:
                    create_content_slide(prs, section_title, bullets)

            # Check if it's a main title (we already handled the main title)
            elif source.startswith("# "):
                continue

            # Check if it contains a table
            elif "|" in source and "---" in source:
                # Extract table title
                table_lines = source.split("\n")
                table_title = ""
                for line in table_lines:
                    if line.startswith("## "):
                        table_title = line[3:]
                    elif line.startswith("### "):
                        table_title = line[4:]
                    elif line.startswith("#### "):
                        table_title = line[5:]

                if not table_title:
                    table_title = "Table Data"

                # Parse table data
                table_data = []
                table_section = False

                for line in table_lines:
                    if "|" in line:
                        if "---" in line:  # Skip the separator line
                            continue

                        cells = [cell.strip() for cell in line.split("|")[1:-1]]
                        if cells:
                            table_data.append(cells)

                if table_data:
                    create_table_slide(prs, table_title, table_data)

            # Check if it contains an image
            elif "![" in source and "](" in source:
                # Extract image information
                img_title = "Image"
                img_path = ""
                img_caption = ""

                # Find title
                for line in source.split("\n"):
                    if line.startswith("### "):
                        img_title = line[4:]
                    elif line.startswith("## "):
                        img_title = line[3:]

                # Extract image path
                img_start = source.find("![")
                if img_start >= 0:
                    img_path_start = source.find("](", img_start) + 2
                    img_path_end = source.find(")", img_path_start)
                    img_path = source[img_path_start:img_path_end]

                    # Check for caption
                    caption_lines = source.split("\n")
                    for i, line in enumerate(caption_lines):
                        if (
                            "![" in line
                            and i + 1 < len(caption_lines)
                            and caption_lines[i + 1].startswith("_")
                        ):
                            img_caption = caption_lines[i + 1].strip("_")

                # For this example, we'll just use a placeholder image path
                # In production, you'd need to download the image or use a local path
                img_placeholder = "placeholder.jpg"
                create_image_slide(prs, img_title, img_placeholder, img_caption)

            # Handle other markdown content (paragraphs, etc.)
            else:
                # Extract a title for the slide
                lines = source.split("\n")
                slide_title = "Information"
                content = []

                for line in lines:
                    if line.startswith("### "):
                        slide_title = line[4:]
                    elif line.startswith("## "):
                        slide_title = line[3:]
                    elif line.strip() and not line.startswith("_") and not line.startswith("!"):
                        content.append(line)

                if content:
                    create_content_slide(prs, slide_title, content)

        elif cell["cell_type"] == "code":
            # Create code slide
            code = "".join(cell["source"])
            if len(code.strip()) > 0:
                create_content_slide(prs, "Code Example", code, is_code=True)

    # Add summary slide
    summary_title = "Summary"
    summary_points = [
        "AI refers to systems that mimic human intelligence to perform tasks",
        "Types of AI: Narrow AI, General AI (AGI), and Superintelligent AI",
        "Real-world applications include virtual assistants, image recognition, recommendation systems, and autonomous vehicles",
        "Most current AI systems are examples of Narrow AI",
        "The field continues to advance rapidly with new breakthroughs",
    ]
    create_content_slide(prs, summary_title, summary_points)

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    notebook_path = r"e:\work\PythonScripts\AI-intro\lesson1_intro_to_ai.ipynb"
    output_path = r"e:\work\PythonScripts\AI-intro\lesson1_intro_to_ai.pptx"
    notebook_to_ppt(notebook_path, output_path)
